#!/usr/bin/env python3
"""
IPR (Individual Performance Report) PowerPoint Generator
Design: AIX R&R tone matching (Navy Blue, Bright Blue, Minimal)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Color Palette (AIX R&R Design)
NAVY_BLUE = RGBColor(35, 38, 105)      # #232669
BRIGHT_BLUE = RGBColor(5, 131, 242)    # #0583F2
LIGHT_GRAY = RGBColor(235, 236, 240)   # #EBECF0
BLACK = RGBColor(0, 0, 0)              # #000000
WHITE = RGBColor(255, 255, 255)        # #FFFFFF
DARK_GRAY = RGBColor(107, 114, 128)    # #6B7280
LIGHT_RED = RGBColor(254, 242, 242)    # #FEF2F2
LIGHT_GREEN = RGBColor(240, 253, 244)  # #F0FDF4
RED = RGBColor(239, 68, 68)            # #EF4444
GREEN = RGBColor(16, 185, 129)         # #10B981
YELLOW = RGBColor(251, 191, 36)        # #FBBF24
ORANGE = RGBColor(249, 115, 22)        # #F97316

def create_presentation():
    """Create the main presentation object"""
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 ratio
    prs.slide_height = Inches(5.625)
    return prs

def add_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background (Light Navy gradient effect)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 246, 250)  # Very light blue-gray
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Individual Performance Report"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    title_para.font.name = 'Montserrat'
    
    # Subtitle (Korean)
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "작업자 성과 관리 시스템"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = NAVY_BLUE
    
    # Subtitle 2
    subtitle2_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.5))
    subtitle2_frame = subtitle2_box.text_frame
    subtitle2_frame.text = "MES R018 Analysis - Long-term Roadmap"
    subtitle2_para = subtitle2_frame.paragraphs[0]
    subtitle2_para.alignment = PP_ALIGN.CENTER
    subtitle2_para.font.size = Pt(18)
    subtitle2_para.font.color.rgb = DARK_GRAY
    
    # Date subtitle
    date_box = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(8), Inches(0.3))
    date_frame = date_box.text_frame
    date_frame.text = "2026 ~ 2028"
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = DARK_GRAY
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(8.5), Inches(5.2), Inches(1), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Feb 2026"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.RIGHT
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = DARK_GRAY

def add_journey_slide(prs):
    """Slide 2: The Journey (AS-IS vs TO-BE)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "The Journey"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    
    # AS-IS Box (Left)
    left_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.2), Inches(4), Inches(3.5)
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = LIGHT_RED
    left_shape.line.color.rgb = RED
    left_shape.line.width = Pt(2)
    
    left_text = left_shape.text_frame
    left_text.text = "AS-IS\n현재의 한계\n\n• 데이터는 있지만 활용 없음\n• 작업자 성과 측정 불가\n• 인력 계획 경험 의존\n• 비효율 지속"
    for para in left_text.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = BLACK
        para.space_after = Pt(8)
    left_text.paragraphs[0].font.size = Pt(20)
    left_text.paragraphs[0].font.bold = True
    
    # TO-BE Box (Right)
    right_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.2), Inches(4), Inches(3.5)
    )
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = LIGHT_GREEN
    right_shape.line.color.rgb = GREEN
    right_shape.line.width = Pt(2)
    
    right_text = right_shape.text_frame
    right_text.text = "TO-BE\n데이터 기반 의사결정\n\n• 3가지 핵심 지표 실시간 집계\n• 개인별 성과 리포트\n• AI 기반 인력 예측\n• 생산성 20% ↑"
    for para in right_text.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = BLACK
        para.space_after = Pt(8)
    right_text.paragraphs[0].font.size = Pt(20)
    right_text.paragraphs[0].font.bold = True

def add_three_metric_slide(prs):
    """Slide 3: 3-Metric System"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "3-Metric System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "작업자 성과의 3가지 차원"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_GRAY
    
    # Card 1: Time Utilization
    card1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.3), Inches(2.8), Inches(2.5)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(243, 244, 246)  # Light purple-gray
    card1.line.color.rgb = BRIGHT_BLUE
    card1.line.width = Pt(1)
    
    card1_text = card1.text_frame
    card1_text.text = "📊 Time Utilization\n작업 시간 활용률\n\n76.3%\n\n(근무 시간 대비\n실제 작업 시간)"
    card1_text.paragraphs[0].font.size = Pt(16)
    card1_text.paragraphs[0].font.bold = True
    card1_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    card1_text.paragraphs[3].font.size = Pt(32)
    card1_text.paragraphs[3].font.bold = True
    card1_text.paragraphs[3].font.color.rgb = BRIGHT_BLUE
    card1_text.paragraphs[3].alignment = PP_ALIGN.CENTER
    
    # Card 2: Work Efficiency
    card2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.8), Inches(1.3), Inches(2.8), Inches(2.5)
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(239, 246, 255)  # Light blue
    card2.line.color.rgb = BRIGHT_BLUE
    card2.line.width = Pt(1)
    
    card2_text = card2.text_frame
    card2_text.text = "⚡ Work Efficiency\n작업 효율성\n\n103.2%\n\n(표준 시간 대비\n실제 생산성)"
    card2_text.paragraphs[0].font.size = Pt(16)
    card2_text.paragraphs[0].font.bold = True
    card2_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    card2_text.paragraphs[3].font.size = Pt(32)
    card2_text.paragraphs[3].font.bold = True
    card2_text.paragraphs[3].font.color.rgb = BRIGHT_BLUE
    card2_text.paragraphs[3].alignment = PP_ALIGN.CENTER
    
    # Card 3: Work Completion
    card3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.3), Inches(2.8), Inches(2.5)
    )
    card3.fill.solid()
    card3.fill.fore_color.rgb = LIGHT_GREEN
    card3.line.color.rgb = GREEN
    card3.line.width = Pt(1)
    
    card3_text = card3.text_frame
    card3_text.text = "✅ Work Completion\n작업 완결성\n\n7.3\njobs/week\n\n(완료된 작업 수량)"
    card3_text.paragraphs[0].font.size = Pt(16)
    card3_text.paragraphs[0].font.bold = True
    card3_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    card3_text.paragraphs[3].font.size = Pt(32)
    card3_text.paragraphs[3].font.bold = True
    card3_text.paragraphs[3].font.color.rgb = GREEN
    card3_text.paragraphs[3].alignment = PP_ALIGN.CENTER
    
    # Bottom summary
    summary_box = slide.shapes.add_textbox(Inches(2), Inches(4.3), Inches(6), Inches(0.5))
    summary_frame = summary_box.text_frame
    summary_frame.text = "→ Overall Score: A- (87.2/100)"
    summary_para = summary_frame.paragraphs[0]
    summary_para.alignment = PP_ALIGN.CENTER
    summary_para.font.size = Pt(20)
    summary_para.font.bold = True
    summary_para.font.color.rgb = NAVY_BLUE

def add_completion_score_slide(prs):
    """Slide 4: Work Completion Score"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Work Completion Score"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "작업 완결성 측정"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_GRAY
    
    # Left: Definition Box
    def_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.3), Inches(4), Inches(3.5)
    )
    def_box.fill.solid()
    def_box.fill.fore_color.rgb = WHITE
    def_box.line.color.rgb = NAVY_BLUE
    def_box.line.width = Pt(2)
    
    def_text = def_box.text_frame
    def_text.text = "📋 Definition\n\nWork Completion Score\n= Σ (Job Rate / 100)\n\nExample:\nW/O-123: 100% → 1.0\nW/O-124: 40% → 0.4\nW/O-125: 80% → 0.8\n\nTotal: 2.2 jobs"
    def_text.paragraphs[0].font.size = Pt(18)
    def_text.paragraphs[0].font.bold = True
    def_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    def_text.paragraphs[2].font.size = Pt(16)
    def_text.paragraphs[2].font.bold = True
    def_text.paragraphs[2].font.color.rgb = BRIGHT_BLUE
    
    # Right: Chart placeholder
    chart_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.3), Inches(4), Inches(3.5)
    )
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = RGBColor(249, 250, 251)
    chart_box.line.color.rgb = LIGHT_GRAY
    chart_box.line.width = Pt(1)
    
    chart_text = chart_box.text_frame
    chart_text.text = "Job Status Breakdown\n\n✅ Completed (100%): 3 jobs\n🟡 Near (80-99%): 2 jobs\n🟠 Progress (50-79%): 1 job\n🔴 Started (<50%): 3 jobs\n\nTotal Score: 5.4 / 9 jobs\nCompletion Rate: 60%"
    chart_text.paragraphs[0].font.size = Pt(18)
    chart_text.paragraphs[0].font.bold = True
    chart_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    chart_text.paragraphs[7].font.size = Pt(16)
    chart_text.paragraphs[7].font.bold = True
    chart_text.paragraphs[7].font.color.rgb = BRIGHT_BLUE

def add_roadmap_timeline_slide(prs):
    """Slide 5: Roadmap Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Roadmap Timeline"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "2026 Q2 ~ 2028 Q2"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_GRAY
    
    # Timeline line
    timeline_line = slide.shapes.add_connector(1, Inches(1), Inches(2.5), Inches(9), Inches(2.5))
    timeline_line.line.color.rgb = BRIGHT_BLUE
    timeline_line.line.width = Pt(3)
    
    # Phase 1
    phase1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.3), Inches(2.5), Inches(2.3)
    )
    phase1.fill.solid()
    phase1.fill.fore_color.rgb = WHITE
    phase1.line.color.rgb = NAVY_BLUE
    phase1.line.width = Pt(2)
    
    p1_text = phase1.text_frame
    p1_text.text = "Phase 1 (0-3 mo)\nFoundation\n2026 Q2 ~ Q3\n\n• 3-Metric System\n• Individual Report\n• Team Dashboard\n\nInvestment: $60K\nROI: 2.5x"
    p1_text.paragraphs[0].font.size = Pt(14)
    p1_text.paragraphs[0].font.bold = True
    p1_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    p1_text.paragraphs[8].font.color.rgb = BRIGHT_BLUE
    p1_text.paragraphs[9].font.color.rgb = GREEN
    p1_text.paragraphs[9].font.bold = True
    
    # Phase 2
    phase2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.6), Inches(1.3), Inches(2.5), Inches(2.3)
    )
    phase2.fill.solid()
    phase2.fill.fore_color.rgb = WHITE
    phase2.line.color.rgb = NAVY_BLUE
    phase2.line.width = Pt(2)
    
    p2_text = phase2.text_frame
    p2_text.text = "Phase 2 (3-6 mo)\nIntelligence\n2026 Q4 ~ 2027 Q2\n\n• S/T Learning\n• MH Forecasting\n• Workforce Planning\n\nInvestment: $120K\nROI: 4.2x"
    p2_text.paragraphs[0].font.size = Pt(14)
    p2_text.paragraphs[0].font.bold = True
    p2_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    p2_text.paragraphs[8].font.color.rgb = BRIGHT_BLUE
    p2_text.paragraphs[9].font.color.rgb = GREEN
    p2_text.paragraphs[9].font.bold = True
    
    # Phase 3
    phase3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.4), Inches(1.3), Inches(2.5), Inches(2.3)
    )
    phase3.fill.solid()
    phase3.fill.fore_color.rgb = WHITE
    phase3.line.color.rgb = NAVY_BLUE
    phase3.line.width = Pt(2)
    
    p3_text = phase3.text_frame
    p3_text.text = "Phase 3 (6-12 mo)\nOptimization\n2027 Q3 ~ 2028 Q2\n\n• AI Task Allocation\n• Cost Analysis\n• ERP Integration\n\nInvestment: $240K\nROI: 5.0x"
    p3_text.paragraphs[0].font.size = Pt(14)
    p3_text.paragraphs[0].font.bold = True
    p3_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    p3_text.paragraphs[8].font.color.rgb = BRIGHT_BLUE
    p3_text.paragraphs[9].font.color.rgb = GREEN
    p3_text.paragraphs[9].font.bold = True
    
    # Bottom summary
    summary_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(7), Inches(0.5))
    summary_frame = summary_box.text_frame
    summary_frame.text = "3-Year Cumulative ROI: 11.1x ($4.68M net gain)"
    summary_para = summary_frame.paragraphs[0]
    summary_para.alignment = PP_ALIGN.CENTER
    summary_para.font.size = Pt(18)
    summary_para.font.bold = True
    summary_para.font.color.rgb = NAVY_BLUE

def add_phase1_slide(prs):
    """Slide 6: Phase 1 - Foundation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Phase 1: Foundation"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "3-Metric System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_GRAY
    
    # Box 1: Raw Data
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.5), Inches(2.3), Inches(2.5)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = WHITE
    box1.line.color.rgb = BRIGHT_BLUE
    box1.line.width = Pt(2)
    
    b1_text = box1.text_frame
    b1_text.text = "📊 Raw Data\n38,280 records\n\n• Work Orders\n• Time Logs\n• Worker Assignments"
    b1_text.paragraphs[0].font.size = Pt(18)
    b1_text.paragraphs[0].font.bold = True
    b1_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    
    # Arrow 1
    arrow1_box = slide.shapes.add_textbox(Inches(3.3), Inches(2.5), Inches(1), Inches(0.5))
    arrow1_frame = arrow1_box.text_frame
    arrow1_frame.text = "→"
    arrow1_para = arrow1_frame.paragraphs[0]
    arrow1_para.alignment = PP_ALIGN.CENTER
    arrow1_para.font.size = Pt(40)
    arrow1_para.font.color.rgb = BRIGHT_BLUE
    
    # Box 2: Engine
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.8), Inches(1.5), Inches(2.3), Inches(2.5)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.color.rgb = BRIGHT_BLUE
    box2.line.width = Pt(2)
    
    b2_text = box2.text_frame
    b2_text.text = "⚙️ 3-Metric Engine\n\nTime Utilization\n+\nWork Efficiency\n+\nWork Completion"
    b2_text.paragraphs[0].font.size = Pt(18)
    b2_text.paragraphs[0].font.bold = True
    b2_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    
    # Arrow 2
    arrow2_box = slide.shapes.add_textbox(Inches(6.3), Inches(2.5), Inches(1), Inches(0.5))
    arrow2_frame = arrow2_box.text_frame
    arrow2_frame.text = "→"
    arrow2_para = arrow2_frame.paragraphs[0]
    arrow2_para.alignment = PP_ALIGN.CENTER
    arrow2_para.font.size = Pt(40)
    arrow2_para.font.color.rgb = BRIGHT_BLUE
    
    # Box 3: Output
    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.5), Inches(2.3), Inches(2.5)
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = WHITE
    box3.line.color.rgb = BRIGHT_BLUE
    box3.line.width = Pt(2)
    
    b3_text = box3.text_frame
    b3_text.text = "📈 Output\n\n• Individual\n  Performance Report\n• Team Dashboard\n• Real-time Score"
    b3_text.paragraphs[0].font.size = Pt(18)
    b3_text.paragraphs[0].font.bold = True
    b3_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    
    # Timeline
    timeline_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(0.5))
    timeline_frame = timeline_box.text_frame
    timeline_frame.text = "M1: Design | M2: Development | M3: Testing | M4: Rollout"
    timeline_para = timeline_frame.paragraphs[0]
    timeline_para.alignment = PP_ALIGN.CENTER
    timeline_para.font.size = Pt(14)
    timeline_para.font.color.rgb = NAVY_BLUE

def add_roi_summary_slide(prs):
    """Slide 9: ROI Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "ROI Summary"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "투자 대비 수익"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_GRAY
    
    # Investment
    inv_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(3), Inches(0.8))
    inv_frame = inv_box.text_frame
    inv_frame.text = "Investment\n$420K"
    inv_para1 = inv_frame.paragraphs[0]
    inv_para1.alignment = PP_ALIGN.CENTER
    inv_para1.font.size = Pt(18)
    inv_para1.font.color.rgb = DARK_GRAY
    inv_para2 = inv_frame.paragraphs[1]
    inv_para2.alignment = PP_ALIGN.CENTER
    inv_para2.font.size = Pt(48)
    inv_para2.font.bold = True
    inv_para2.font.color.rgb = RED
    
    # Arrow
    arrow_box = slide.shapes.add_textbox(Inches(4), Inches(2.2), Inches(1.5), Inches(0.5))
    arrow_frame = arrow_box.text_frame
    arrow_frame.text = "→"
    arrow_para = arrow_frame.paragraphs[0]
    arrow_para.alignment = PP_ALIGN.CENTER
    arrow_para.font.size = Pt(60)
    arrow_para.font.color.rgb = BRIGHT_BLUE
    
    # Returns
    ret_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(3), Inches(0.8))
    ret_frame = ret_box.text_frame
    ret_frame.text = "Returns\n$1.85M"
    ret_para1 = ret_frame.paragraphs[0]
    ret_para1.alignment = PP_ALIGN.CENTER
    ret_para1.font.size = Pt(18)
    ret_para1.font.color.rgb = DARK_GRAY
    ret_para2 = ret_frame.paragraphs[1]
    ret_para2.alignment = PP_ALIGN.CENTER
    ret_para2.font.size = Pt(48)
    ret_para2.font.bold = True
    ret_para2.font.color.rgb = GREEN
    
    # Net Gain
    net_box = slide.shapes.add_textbox(Inches(2.5), Inches(3), Inches(5), Inches(0.6))
    net_frame = net_box.text_frame
    net_frame.text = "Net Gain: $1.43M"
    net_para = net_frame.paragraphs[0]
    net_para.alignment = PP_ALIGN.CENTER
    net_para.font.size = Pt(28)
    net_para.font.bold = True
    net_para.font.color.rgb = BRIGHT_BLUE
    
    # ROI Box (Highlighted)
    roi_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(3.8), Inches(4), Inches(1)
    )
    roi_shape.fill.solid()
    roi_shape.fill.fore_color.rgb = RGBColor(254, 243, 199)  # Yellow highlight
    roi_shape.line.color.rgb = NAVY_BLUE
    roi_shape.line.width = Pt(3)
    
    roi_text = roi_shape.text_frame
    roi_text.text = "ROI: 11.1x"
    roi_para = roi_text.paragraphs[0]
    roi_para.alignment = PP_ALIGN.CENTER
    roi_para.font.size = Pt(60)
    roi_para.font.bold = True
    roi_para.font.color.rgb = NAVY_BLUE
    roi_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Payback Period
    payback_box = slide.shapes.add_textbox(Inches(2.5), Inches(5), Inches(5), Inches(0.4))
    payback_frame = payback_box.text_frame
    payback_frame.text = "Payback Period: 9 months"
    payback_para = payback_frame.paragraphs[0]
    payback_para.alignment = PP_ALIGN.CENTER
    payback_para.font.size = Pt(18)
    payback_para.font.bold = True
    payback_para.font.color.rgb = NAVY_BLUE

def add_vision_2028_slide(prs):
    """Slide 11: Vision 2028"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Vision 2028"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Manufacturing Intelligence"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_GRAY
    
    # Center: AI Copilot
    center_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(5), Inches(0.8))
    center_frame = center_box.text_frame
    center_frame.text = "AI Copilot\n작업 현장의 AI 비서"
    center_para1 = center_frame.paragraphs[0]
    center_para1.alignment = PP_ALIGN.CENTER
    center_para1.font.size = Pt(40)
    center_para1.font.bold = True
    center_para1.font.color.rgb = NAVY_BLUE
    center_para2 = center_frame.paragraphs[1]
    center_para2.alignment = PP_ALIGN.CENTER
    center_para2.font.size = Pt(18)
    center_para2.font.color.rgb = DARK_GRAY
    
    # 4 Icons in circular layout
    # Icon 1: Top
    icon1 = slide.shapes.add_textbox(Inches(4), Inches(2.4), Inches(2), Inches(0.8))
    icon1_frame = icon1.text_frame
    icon1_frame.text = "🔮 Predict\n주간 생산 계획\n자동 생성"
    icon1_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon1_frame.paragraphs[0].font.size = Pt(16)
    icon1_frame.paragraphs[0].font.bold = True
    icon1_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    
    # Icon 2: Right
    icon2 = slide.shapes.add_textbox(Inches(6.5), Inches(3.2), Inches(2), Inches(0.8))
    icon2_frame = icon2.text_frame
    icon2_frame.text = "⚙️ Optimize\n실시간 자원\n최적 배분"
    icon2_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon2_frame.paragraphs[0].font.size = Pt(16)
    icon2_frame.paragraphs[0].font.bold = True
    icon2_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    
    # Icon 3: Bottom
    icon3 = slide.shapes.add_textbox(Inches(4), Inches(4), Inches(2), Inches(0.8))
    icon3_frame = icon3.text_frame
    icon3_frame.text = "🚨 Alert\n이상 징후\n사전 경보"
    icon3_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon3_frame.paragraphs[0].font.size = Pt(16)
    icon3_frame.paragraphs[0].font.bold = True
    icon3_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    
    # Icon 4: Left
    icon4 = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(2), Inches(0.8))
    icon4_frame = icon4.text_frame
    icon4_frame.text = "📈 Learn\n지속적 성능\n개선"
    icon4_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon4_frame.paragraphs[0].font.size = Pt(16)
    icon4_frame.paragraphs[0].font.bold = True
    icon4_frame.paragraphs[0].font.color.rgb = NAVY_BLUE
    
    # Bottom
    bottom_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(0.4))
    bottom_frame = bottom_box.text_frame
    bottom_frame.text = "→ Autonomous Manufacturing Operations"
    bottom_para = bottom_frame.paragraphs[0]
    bottom_para.alignment = PP_ALIGN.CENTER
    bottom_para.font.size = Pt(20)
    bottom_para.font.bold = True
    bottom_para.font.color.rgb = NAVY_BLUE

def add_call_to_action_slide(prs):
    """Slide 12: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Next Steps"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY_BLUE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.3))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "지금 시작하세요"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = DARK_GRAY
    
    # Checklist
    checklist_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.3), Inches(7), Inches(2))
    checklist_frame = checklist_box.text_frame
    checklist_frame.text = "✅ Phase 1 Budget Approval - $60K (2026 Q2)\n\n✅ Work Completion Score Launch - 4 weeks\n\n✅ IPR Pilot Program - 50 workers, 8 weeks\n\n✅ Team Dashboard Rollout - 12 weeks"
    for para in checklist_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = BLACK
        para.space_after = Pt(6)
    
    # Timeline
    timeline_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(0.5))
    timeline_frame = timeline_box.text_frame
    timeline_frame.text = "M1 → M2 → M3 → M4\nDesign | Develop | Test | Deploy"
    timeline_para = timeline_frame.paragraphs[0]
    timeline_para.alignment = PP_ALIGN.CENTER
    timeline_para.font.size = Pt(14)
    timeline_para.font.bold = True
    timeline_para.font.color.rgb = BRIGHT_BLUE
    
    # CTA Box
    cta_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(4.2), Inches(6), Inches(0.8)
    )
    cta_shape.fill.solid()
    cta_shape.fill.fore_color.rgb = RGBColor(254, 243, 199)  # Yellow
    cta_shape.line.color.rgb = NAVY_BLUE
    cta_shape.line.width = Pt(3)
    
    cta_text = cta_shape.text_frame
    cta_text.text = "🚀 Start Today\nExecutive Sponsorship Required | Change Management Plan | IT Team Alignment"
    cta_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    cta_text.paragraphs[0].font.size = Pt(20)
    cta_text.paragraphs[0].font.bold = True
    cta_text.paragraphs[0].font.color.rgb = NAVY_BLUE
    cta_text.paragraphs[1].alignment = PP_ALIGN.CENTER
    cta_text.paragraphs[1].font.size = Pt(12)
    cta_text.paragraphs[1].font.color.rgb = BLACK
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(3), Inches(5.2), Inches(4), Inches(0.3))
    contact_frame = contact_box.text_frame
    contact_frame.text = "MES Team | mes@company.com"
    contact_para = contact_frame.paragraphs[0]
    contact_para.alignment = PP_ALIGN.CENTER
    contact_para.font.size = Pt(12)
    contact_para.font.color.rgb = DARK_GRAY

def main():
    """Main function to generate the presentation"""
    print("🚀 Starting IPR PowerPoint Generation...")
    
    # Create presentation
    prs = create_presentation()
    
    # Add all slides
    print("📄 Adding Slide 1: Title")
    add_title_slide(prs)
    
    print("📄 Adding Slide 2: The Journey")
    add_journey_slide(prs)
    
    print("📄 Adding Slide 3: 3-Metric System")
    add_three_metric_slide(prs)
    
    print("📄 Adding Slide 4: Work Completion Score")
    add_completion_score_slide(prs)
    
    print("📄 Adding Slide 5: Roadmap Timeline")
    add_roadmap_timeline_slide(prs)
    
    print("📄 Adding Slide 6: Phase 1 - Foundation")
    add_phase1_slide(prs)
    
    print("📄 Adding Slide 9: ROI Summary")
    add_roi_summary_slide(prs)
    
    print("📄 Adding Slide 11: Vision 2028")
    add_vision_2028_slide(prs)
    
    print("📄 Adding Slide 12: Call to Action")
    add_call_to_action_slide(prs)
    
    # Save presentation
    output_path = '/home/user/webapp/IPR_Roadmap_Presentation.pptx'
    prs.save(output_path)
    
    print(f"\n✅ Presentation saved: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\n🎉 PowerPoint generation complete!")

if __name__ == '__main__':
    main()
